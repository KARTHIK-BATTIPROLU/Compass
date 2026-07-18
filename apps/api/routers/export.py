import io
import csv
import json
import re
from fastapi import APIRouter, HTTPException, Depends
from fastapi.responses import StreamingResponse
from pptx import Presentation
import docx
from agent.auth import get_supabase, get_current_user, user_owns_session

router = APIRouter(prefix="/api/artifacts", tags=["export"])

@router.get("/{artifact_id}/export")
async def export_artifact(
    artifact_id: str,
    format: str = "docx",
    user = Depends(get_current_user)
):
    sb = get_supabase()
    res = sb.table("artifacts").select("*").eq("id", artifact_id).execute()
    if not res.data:
        raise HTTPException(status_code=404, detail="Artifact not found")
        
    artifact = res.data[0]
    session_id = artifact.get("session_id")
    
    if session_id and not user_owns_session(user.id, session_id):
        raise HTTPException(status_code=403, detail="Forbidden")

    content_md = artifact.get("content_md", "")
    title = f"Export_{artifact_id}"

    if format == "pptx":
        return export_pptx(content_md, title)
    elif format == "csv":
        return export_csv(content_md, title)
    elif format == "pdf":
        # Fallback to DOCX on systems without PDF engines
        return export_docx(content_md, title, is_pdf_fallback=True)
    else:
        return export_docx(content_md, title)

def export_pptx(md: str, title: str):
    prs = Presentation()
    # Simple split by horizontal rules or headers
    slides_content = [s for s in md.split("---") if s.strip()]
    
    for slide_text in slides_content:
        lines = [line.strip() for line in slide_text.strip().split("\n") if line.strip()]
        if not lines: continue
        
        slide = prs.slides.add_slide(prs.slide_layouts[1]) # Title and Content
        slide.shapes.title.text = lines[0].replace("#", "").strip()
        
        if len(lines) > 1:
            tf = slide.placeholders[1].text_frame
            tf.text = "\n".join(lines[1:])
            
    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return StreamingResponse(
        bio,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={title}.pptx"}
    )

def export_docx(md: str, title: str, is_pdf_fallback=False):
    doc = docx.Document()
    doc.add_heading(title, 0)
    
    if is_pdf_fallback:
        doc.add_paragraph("Note: PDF export is not supported natively in this environment. Falling back to DOCX format.")
        
    for para in md.split("\n\n"):
        if para.strip():
            doc.add_paragraph(para.strip())
            
    bio = io.BytesIO()
    doc.save(bio)
    bio.seek(0)
    ext = "docx" if not is_pdf_fallback else "pdf-fallback.docx"
    return StreamingResponse(
        bio,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f"attachment; filename={title}.{ext}"}
    )

def export_csv(md: str, title: str):
    # Flashcards are stored as JSON (title/cards[]) inside an <artifact> tag.
    bio = io.StringIO()
    writer = csv.writer(bio)
    writer.writerow(["Front", "Back"])

    cards = []
    match = re.search(r"<artifact[^>]*>(.*?)</artifact>", md, re.DOTALL)
    raw = match.group(1).strip() if match else md.strip()
    try:
        data = json.loads(raw)
        cards = data.get("cards", [])
    except Exception:
        cards = []

    if cards:
        for c in cards:
            writer.writerow([c.get("front", ""), c.get("back", "")])
    else:
        # Fallback for non-JSON content: naive "Term: definition" line parsing
        for line in md.split("\n"):
            if ":" in line:
                parts = line.split(":", 1)
                writer.writerow([parts[0].strip(), parts[1].strip()])

    bio.seek(0)
    return StreamingResponse(
        iter([bio.getvalue()]),
        media_type="text/csv",
        headers={"Content-Disposition": f"attachment; filename={title}.csv"}
    )
